import os
import uuid
import shutil
import tempfile
import subprocess
from pathlib import Path

import fitz

from PIL import Image, ImageChops

from docx import Document
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse

from pdf2docx import Converter

# ==========================================================
# QUICKFILE TOOLS BACKEND
# ==========================================================

app = FastAPI(title="QuickFile Tools Backend")


# ==========================================================
# CORS
# ==========================================================

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ==========================================================
# TEMP FOLDER
# Project folder ke bahar temp files rakhta hai
# Live Server auto-refresh problem ko rokta hai
# ==========================================================

TEMP_DIR = Path(tempfile.gettempdir()) / "quickfile_tools_temp"

TEMP_DIR.mkdir(parents=True, exist_ok=True)


# ==========================================================
# HOME TEST
# ==========================================================


@app.get("/")
def home():

    return {"status": "ok", "message": "QuickFile Tools backend is running"}


# ==========================================================
# PDF TO WORD - EDITABLE MODE
# pdf2docx use karta hai
# ==========================================================


@app.post("/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):

    original_name = file.filename or "document.pdf"

    if not original_name.lower().endswith(".pdf"):

        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    job_id = str(uuid.uuid4())

    pdf_path = TEMP_DIR / f"{job_id}.pdf"

    docx_path = TEMP_DIR / f"{job_id}.docx"

    try:

        # Save uploaded PDF
        with open(pdf_path, "wb") as buffer:

            shutil.copyfileobj(file.file, buffer)

        # Convert PDF to DOCX
        converter = Converter(str(pdf_path))

        try:

            converter.convert(str(docx_path), start=0, end=None)

        finally:

            converter.close()

        if not docx_path.exists():

            raise HTTPException(
                status_code=500, detail=("Word file could not be created.")
            )

        download_name = Path(original_name).stem + ".docx"

        return FileResponse(
            path=str(docx_path),
            media_type=(
                "application/"
                "vnd.openxmlformats-officedocument."
                "wordprocessingml.document"
            ),
            filename=download_name,
        )

    except HTTPException:

        raise

    except Exception as error:

        print("PDF to Word Error:", error)

        raise HTTPException(status_code=500, detail=("PDF to Word conversion failed."))

    finally:

        if pdf_path.exists():

            try:
                os.remove(pdf_path)

            except Exception:
                pass


# ==========================================================
# PDF TO WORD - PRESERVE ORIGINAL LAYOUT
# PDF page ko image me convert karke
# A4 Word page ke center me rakhta hai
# ==========================================================


@app.post("/pdf-to-word-layout")
async def pdf_to_word_layout(file: UploadFile = File(...)):

    original_name = file.filename or "document.pdf"

    if not original_name.lower().endswith(".pdf"):

        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    job_id = str(uuid.uuid4())

    pdf_path = TEMP_DIR / f"{job_id}.pdf"

    docx_path = TEMP_DIR / f"{job_id}-layout.docx"

    image_paths = []

    pdf_document = None

    try:

        # Save uploaded PDF
        with open(pdf_path, "wb") as buffer:

            shutil.copyfileobj(file.file, buffer)

        # Open PDF
        pdf_document = fitz.open(str(pdf_path))

        # Create Word document
        document = Document()

        section = document.sections[0]

        # A4 page size
        section.page_width = Inches(8.27)
        section.page_height = Inches(11.69)

        # Equal margins
        section.top_margin = Inches(0.30)
        section.bottom_margin = Inches(0.30)
        section.left_margin = Inches(0.30)
        section.right_margin = Inches(0.30)

        # Convert every PDF page
        for index, page in enumerate(pdf_document):

            # High quality render
            matrix = fitz.Matrix(2, 2)

            pix = page.get_pixmap(matrix=matrix, alpha=False)

            image_path = TEMP_DIR / f"{job_id}_{index}.png"

            pix.save(str(image_path))

            image_paths.append(image_path)

            # ----------------------------------------------
            # Remove unnecessary white space
            # ----------------------------------------------

            img = Image.open(image_path).convert("RGB")

            white_background = Image.new("RGB", img.size, (255, 255, 255))

            difference = ImageChops.difference(img, white_background)

            bbox = difference.getbbox()

            if bbox:

                img = img.crop(bbox)

                img.save(image_path)

            # ----------------------------------------------
            # Center PDF content on Word A4 page
            # ----------------------------------------------

            paragraph = document.add_paragraph()

            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

            paragraph.paragraph_format.left_indent = Inches(0)

            paragraph.paragraph_format.right_indent = Inches(0)

            paragraph.paragraph_format.space_before = 0
            paragraph.paragraph_format.space_after = 0

            run = paragraph.add_run()

            image_width = img.width
            image_height = img.height

            aspect_ratio = image_width / image_height

            # User-tested centered width
            max_width = 6.5

            max_height = 11.09

            if (max_width / aspect_ratio) <= max_height:

                picture_width = max_width

                picture_height = max_width / aspect_ratio

            else:

                picture_height = max_height

                picture_width = max_height * aspect_ratio

            run.add_picture(
                str(image_path),
                width=Inches(picture_width),
                height=Inches(picture_height),
            )

            # Add page break
            if index < (len(pdf_document) - 1):

                document.add_page_break()

        pdf_document.close()

        pdf_document = None

        # Save Word document
        document.save(str(docx_path))

        download_name = Path(original_name).stem + "-layout.docx"

        return FileResponse(
            path=str(docx_path),
            media_type=(
                "application/"
                "vnd.openxmlformats-officedocument."
                "wordprocessingml.document"
            ),
            filename=download_name,
        )

    except Exception as error:

        print("PDF to Word Layout Error:", error)

        raise HTTPException(status_code=500, detail=("Layout conversion failed."))

    finally:

        if pdf_document is not None:

            try:
                pdf_document.close()

            except Exception:
                pass

        if pdf_path.exists():

            try:
                os.remove(pdf_path)

            except Exception:
                pass

        for image_path in image_paths:

            if image_path.exists():

                try:
                    os.remove(image_path)

                except Exception:
                    pass


# ==========================================================
# PDF TO WORD - EDITABLE TABLE / TEXT MODE
# Tables milne par editable Word table banata hai
# ==========================================================


@app.post("/pdf-to-word-ocr")
async def pdf_to_word_ocr(file: UploadFile = File(...)):

    original_name = file.filename or "document.pdf"

    if not original_name.lower().endswith(".pdf"):

        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    job_id = str(uuid.uuid4())

    pdf_path = TEMP_DIR / f"{job_id}.pdf"

    docx_path = TEMP_DIR / f"{job_id}-ocr.docx"

    pdf_document = None

    try:

        # Save PDF
        with open(pdf_path, "wb") as buffer:

            shutil.copyfileobj(file.file, buffer)

        pdf_document = fitz.open(str(pdf_path))

        document = Document()

        section = document.sections[0]

        section.page_width = Inches(8.27)
        section.page_height = Inches(11.69)

        section.top_margin = Inches(0.40)
        section.bottom_margin = Inches(0.40)
        section.left_margin = Inches(0.40)
        section.right_margin = Inches(0.40)

        for page_number, page in enumerate(pdf_document, start=1):

            tables_found = False

            # ----------------------------------------------
            # Try detecting tables
            # ----------------------------------------------

            try:

                tables = page.find_tables()

                if tables and tables.tables:

                    tables_found = True

                    for pdf_table in tables.tables:

                        data = pdf_table.extract()

                        if not data:
                            continue

                        rows = len(data)

                        cols = max(len(row) for row in data if row)
                        word_table = document.add_table(rows=rows, cols=cols)
                        word_table.style = "Table Grid"
                        word_table.alignment = WD_TABLE_ALIGNMENT.CENTER
                        word_table.autofit = False

                        section = document.sections[0]

                        available_width = (
                            section.page_width
                            - section.left_margin
                            - section.right_margin
                        )

                        column_width = int(available_width / cols)

                        for column in word_table.columns:
                            column.width = column_width

                        for row in word_table.rows:
                            for cell in row.cells:
                                cell.width = column_width

                        for row_index, row_data in enumerate(data):

                            for col_index in range(cols):

                                cell_text = ""

                                if row_data and col_index < len(row_data):

                                    value = row_data[col_index]

                                    if value is not None:

                                        cell_text = str(value)

                                        cell = word_table.cell(row_index, col_index)
                                        cell.text = cell_text

                                        document.add_paragraph()

            except Exception as table_error:

                print("Table detection warning:", table_error)

            # ----------------------------------------------
            # If no table, extract normal text
            # ----------------------------------------------

            if not tables_found:

                text = page.get_text("text").strip()

                if text:

                    document.add_paragraph(text)

                else:

                    document.add_paragraph(
                        (f"[Page {page_number}: " "No editable text detected]")
                    )

            if page_number < len(pdf_document):

                document.add_page_break()

        pdf_document.close()

        pdf_document = None

        document.save(str(docx_path))

        download_name = Path(original_name).stem + "-ocr.docx"

        return FileResponse(
            path=str(docx_path),
            media_type=(
                "application/"
                "vnd.openxmlformats-officedocument."
                "wordprocessingml.document"
            ),
            filename=download_name,
        )

    except Exception as error:

        print("PDF to Word OCR Error:", error)

        raise HTTPException(status_code=500, detail=("OCR conversion failed."))

    finally:

        if pdf_document is not None:

            try:
                pdf_document.close()

            except Exception:
                pass

        if pdf_path.exists():

            try:
                os.remove(pdf_path)

            except Exception:
                pass


# ==========================================================
# WORD TO PDF - LIBREOFFICE
# ==========================================================


@app.post("/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):

    original_name = file.filename or "document.docx"

    if not original_name.lower().endswith((".doc", ".docx")):
        raise HTTPException(status_code=400, detail="Please upload a Word file.")

    job_id = str(uuid.uuid4())

    input_path = TEMP_DIR / f"{job_id}-{original_name}"
    output_dir = TEMP_DIR / f"{job_id}-output"

    output_dir.mkdir(parents=True, exist_ok=True)

    try:
        # Save uploaded Word file
        with open(input_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        soffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"

        command = [
            soffice_path,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(input_path),
        ]

        result = subprocess.run(
            command,
            capture_output=True,
            text=True,
            timeout=120,
        )

        print("LibreOffice output:", result.stdout)
        print("LibreOffice error:", result.stderr)

        pdf_files = list(output_dir.glob("*.pdf"))

        if not pdf_files:
            raise HTTPException(
                status_code=500, detail="PDF file could not be created."
            )

        pdf_path = pdf_files[0]

        download_name = Path(original_name).stem + ".pdf"

        return FileResponse(
            path=str(pdf_path),
            media_type="application/pdf",
            filename=download_name,
        )

    except HTTPException:
        raise

    except Exception as error:
        print("Word to PDF Error:", error)

        raise HTTPException(status_code=500, detail="Word to PDF conversion failed.")

    finally:
        if input_path.exists():
            try:
                os.remove(input_path)
            except Exception:
                pass


# ==========================================================
# PDF TO JPG - PYMUPDF
# ==========================================================


@app.post("/pdf-to-jpg")
async def pdf_to_jpg(file: UploadFile = File(...)):

    original_name = file.filename or "document.pdf"

    if not original_name.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    job_id = str(uuid.uuid4())

    pdf_path = TEMP_DIR / f"{job_id}.pdf"
    output_dir = TEMP_DIR / f"{job_id}-jpg"

    output_dir.mkdir(parents=True, exist_ok=True)

    pdf_document = None

    try:
        with open(pdf_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        pdf_document = fitz.open(str(pdf_path))

        jpg_files = []

        for index, page in enumerate(pdf_document, start=1):

            matrix = fitz.Matrix(2.5, 2.5)

            pix = page.get_pixmap(matrix=matrix, alpha=False)

            jpg_path = output_dir / f"page-{index}.jpg"

            pix.save(str(jpg_path))

            jpg_files.append(jpg_path)

        if not jpg_files:
            raise HTTPException(status_code=500, detail="No JPG files were created.")

        if len(jpg_files) == 1:
            return FileResponse(
                path=str(jpg_files[0]),
                media_type="image/jpeg",
                filename=Path(original_name).stem + ".jpg",
            )

        zip_path = TEMP_DIR / f"{job_id}-jpg.zip"

        shutil.make_archive(str(zip_path.with_suffix("")), "zip", str(output_dir))

        return FileResponse(
            path=str(zip_path),
            media_type="application/zip",
            filename=Path(original_name).stem + "-jpg.zip",
        )

    except HTTPException:
        raise

    except Exception as error:
        print("PDF to JPG Error:", error)

        raise HTTPException(status_code=500, detail="PDF to JPG conversion failed.")

    finally:
        if pdf_document is not None:
            try:
                pdf_document.close()
            except Exception:
                pass

        if pdf_path.exists():
            try:
                os.remove(pdf_path)
            except Exception:
                pass


import os
import tempfile

import fitz
from fastapi import File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse


@app.post("/unlock-pdf")
async def unlock_pdf(file: UploadFile = File(...), password: str = Form("")):
    input_path = None
    output_path = None

    try:
        if not file.filename.lower().endswith(".pdf"):
            raise HTTPException(status_code=400, detail="Please select a PDF file.")

        pdf_bytes = await file.read()

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_input:
            temp_input.write(pdf_bytes)
            input_path = temp_input.name

        output_fd, output_path = tempfile.mkstemp(suffix=".pdf")
        os.close(output_fd)

        document = fitz.open(input_path)

        if document.needs_pass:
            if not password:
                document.close()
                raise HTTPException(
                    status_code=400, detail="This PDF requires a password."
                )

            authenticated = document.authenticate(password)

            if not authenticated:
                document.close()
                raise HTTPException(status_code=401, detail="Incorrect PDF password.")

        document.save(
            output_path, encryption=fitz.PDF_ENCRYPT_NONE, garbage=4, deflate=True
        )

        document.close()

        original_name = os.path.splitext(file.filename)[0]

        return FileResponse(
            output_path,
            media_type="application/pdf",
            filename=f"{original_name}-unlocked.pdf",
        )

    except HTTPException:
        raise

    except Exception as error:
        print("Unlock PDF Error:", error)

        raise HTTPException(status_code=500, detail="Unable to unlock this PDF.")

    finally:
        if input_path and os.path.exists(input_path):
            try:
                os.remove(input_path)
            except Exception:
                pass


@app.post("/protect-pdf")
async def protect_pdf(
    file: UploadFile = File(...),
    password: str = Form(...),
):
    input_path = None
    output_path = None

    try:
        if not file.filename.lower().endswith(".pdf"):
            raise HTTPException(status_code=400, detail="Please select a PDF file.")

        password = password.strip()

        if len(password) < 4:
            raise HTTPException(
                status_code=400, detail="Password must be at least 4 characters."
            )

        pdf_bytes = await file.read()

        if not pdf_bytes:
            raise HTTPException(status_code=400, detail="The selected PDF is empty.")

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_input:
            temp_input.write(pdf_bytes)
            input_path = temp_input.name

        output_fd, output_path = tempfile.mkstemp(suffix=".pdf")
        os.close(output_fd)

        document = fitz.open(input_path)

        if document.needs_pass:
            document.close()

            raise HTTPException(
                status_code=400, detail="This PDF is already password protected."
            )

        document.save(
            output_path,
            encryption=fitz.PDF_ENCRYPT_AES_256,
            user_pw=password,
            owner_pw=password,
            garbage=4,
            deflate=True,
        )

        document.close()

        original_name = os.path.splitext(file.filename)[0]

        return FileResponse(
            output_path,
            media_type="application/pdf",
            filename=f"{original_name}-protected.pdf",
        )

    except HTTPException:
        raise

    except Exception as error:
        print("Protect PDF Error:", error)

        raise HTTPException(status_code=500, detail="Unable to protect this PDF.")

    finally:
        if input_path and os.path.exists(input_path):
            try:
                os.remove(input_path)
            except Exception:
                pass


from io import BytesIO

import pymupdf
from fastapi import File, HTTPException, UploadFile
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.util import Inches


@app.post("/pdf-to-powerpoint")
async def pdf_to_powerpoint(file: UploadFile = File(...)):
    try:
        if not file.filename:
            raise HTTPException(status_code=400, detail="Please select a PDF file.")

        if not file.filename.lower().endswith(".pdf"):
            raise HTTPException(
                status_code=400, detail="Please select a valid PDF file."
            )

        pdf_bytes = await file.read()

        if not pdf_bytes:
            raise HTTPException(
                status_code=400, detail="The selected PDF file is empty."
            )

        try:
            document = pymupdf.open(stream=pdf_bytes, filetype="pdf")
        except Exception:
            raise HTTPException(status_code=400, detail="Unable to open this PDF file.")

        if document.needs_pass:
            document.close()

            raise HTTPException(
                status_code=400,
                detail=("This PDF is password protected. " "Please unlock it first."),
            )

        if document.page_count == 0:
            document.close()

            raise HTTPException(
                status_code=400, detail="This PDF does not contain any pages."
            )

        presentation = Presentation()

        # Remove default first slide if present.
        while len(presentation.slides) > 0:
            slide_id = presentation.slides._sldIdLst[0]
            presentation.part.drop_rel(slide_id.rId)
            presentation.slides._sldIdLst.remove(slide_id)

        first_page = document.load_page(0)
        first_rect = first_page.rect

        if first_rect.width <= 0 or first_rect.height <= 0:
            document.close()

            raise HTTPException(status_code=400, detail="Invalid PDF page dimensions.")

        # Set PowerPoint slide size according to
        # the aspect ratio of the first PDF page.
        slide_width_inches = 10.0

        slide_height_inches = slide_width_inches * first_rect.height / first_rect.width

        # Keep slide dimensions inside a practical range.
        if slide_height_inches > 13.0:
            slide_height_inches = 13.0
            slide_width_inches = (
                slide_height_inches * first_rect.width / first_rect.height
            )

        if slide_height_inches < 5.0:
            slide_height_inches = 5.0
            slide_width_inches = (
                slide_height_inches * first_rect.width / first_rect.height
            )

        presentation.slide_width = Inches(slide_width_inches)

        presentation.slide_height = Inches(slide_height_inches)

        blank_layout = presentation.slide_layouts[6]

        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        # Render each PDF page at approximately
        # 144 DPI for good PowerPoint quality.
        zoom = 2.0
        matrix = pymupdf.Matrix(zoom, zoom)

        for page_number in range(document.page_count):
            page = document.load_page(page_number)

            page_rect = page.rect

            pixmap = page.get_pixmap(matrix=matrix, alpha=False)

            image_bytes = pixmap.tobytes("png")

            image_stream = BytesIO(image_bytes)

            slide = presentation.slides.add_slide(blank_layout)

            page_ratio = page_rect.width / page_rect.height

            slide_ratio = slide_width / slide_height

            if page_ratio > slide_ratio:
                picture_width = slide_width

                picture_height = int(picture_width / page_ratio)

                left = 0

                top = int((slide_height - picture_height) / 2)

            else:
                picture_height = slide_height

                picture_width = int(picture_height * page_ratio)

                top = 0

                left = int((slide_width - picture_width) / 2)

            slide.shapes.add_picture(
                image_stream, left, top, width=picture_width, height=picture_height
            )

            image_stream.close()

        document.close()

        output_stream = BytesIO()

        presentation.save(output_stream)

        output_stream.seek(0)

        original_name = file.filename.rsplit(".", 1)[0]

        safe_name = "".join(
            character if (character.isalnum() or character in " -_()") else "_"
            for character in original_name
        ).strip()

        if not safe_name:
            safe_name = "converted"

        headers = {"Content-Disposition": f'attachment; filename="{safe_name}.pptx"'}

        return StreamingResponse(
            output_stream,
            media_type=(
                "application/vnd.openxmlformats-"
                "officedocument.presentationml.presentation"
            ),
            headers=headers,
        )

    except HTTPException:
        raise

    except Exception as error:
        print("PDF to PowerPoint Error:", error)

        raise HTTPException(
            status_code=500, detail="Unable to convert this PDF to PowerPoint."
        )

    from io import BytesIO


import os

import pymupdf
import pytesseract

from PIL import Image

from fastapi import File, HTTPException, UploadFile
from fastapi.responses import StreamingResponse

from openpyxl import Workbook
from openpyxl.styles import Alignment, Font
from openpyxl.utils import get_column_letter

# --------------------------------------------------
# TESSERACT OCR PATH
# --------------------------------------------------

possible_tesseract_paths = [
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
]

for tesseract_path in possible_tesseract_paths:
    if os.path.exists(tesseract_path):

        pytesseract.pytesseract.tesseract_cmd = tesseract_path

        break


# --------------------------------------------------
# OCR HELPER
# --------------------------------------------------


def extract_ocr_rows(page):

    try:

        pixmap = page.get_pixmap(matrix=pymupdf.Matrix(2.5, 2.5), alpha=False)

        image = Image.open(BytesIO(pixmap.tobytes("png")))

        ocr_data = pytesseract.image_to_data(
            image, output_type=pytesseract.Output.DICT, config="--psm 6"
        )

        grouped_lines = {}

        total_items = len(ocr_data["text"])

        for index in range(total_items):

            text = (ocr_data["text"][index] or "").strip()

            if not text:
                continue

            try:
                confidence = float(ocr_data["conf"][index])
            except Exception:
                confidence = -1

            if confidence < 20:
                continue

            block_number = ocr_data["block_num"][index]

            paragraph_number = ocr_data["par_num"][index]

            line_number = ocr_data["line_num"][index]

            line_key = (block_number, paragraph_number, line_number)

            word_data = {
                "text": text,
                "left": ocr_data["left"][index],
                "top": ocr_data["top"][index],
                "width": ocr_data["width"][index],
            }

            grouped_lines.setdefault(line_key, []).append(word_data)

        rows = []

        for line_words in grouped_lines.values():

            line_words.sort(key=lambda item: item["left"])

            if not line_words:
                continue

            row = []

            current_column = ""

            previous_right = None

            average_height = max([word.get("width", 20) for word in line_words] or [20])

            gap_threshold = max(35, int(average_height * 0.7))

            for word in line_words:

                left = word["left"]

                right = word["left"] + word["width"]

                if previous_right is None:

                    current_column = word["text"]

                else:

                    gap = left - previous_right

                    if gap > gap_threshold:

                        if current_column.strip():

                            row.append(current_column.strip())

                        current_column = word["text"]

                    else:

                        current_column += " " + word["text"]

                previous_right = right

            if current_column.strip():

                row.append(current_column.strip())

            if row:
                rows.append(row)

        return rows

    except Exception as error:

        print("OCR extraction error:", error)

        return []


# --------------------------------------------------
# AUTO COLUMN WIDTH
# --------------------------------------------------


def format_excel_sheet(worksheet):

    for row in worksheet.iter_rows():

        for cell in row:

            cell.alignment = Alignment(vertical="top", wrap_text=True)

    for column_cells in worksheet.columns:

        max_length = 0

        column_number = column_cells[0].column

        for cell in column_cells:

            try:

                cell_value = str(cell.value if cell.value is not None else "")

                if len(cell_value) > max_length:

                    max_length = len(cell_value)

            except Exception:
                pass

        width = min(max(max_length + 2, 12), 50)

        worksheet.column_dimensions[get_column_letter(column_number)].width = width


# --------------------------------------------------
# PDF TO EXCEL - TEXT + OCR + IMAGES
# --------------------------------------------------

from io import BytesIO
import os

import pymupdf
import pytesseract
from PIL import Image as PILImage
from fastapi import File, HTTPException, UploadFile
from fastapi.responses import StreamingResponse
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage
from openpyxl.styles import Alignment, Font
from openpyxl.utils import get_column_letter

# --------------------------------------------------
# TESSERACT PATH
# --------------------------------------------------

possible_tesseract_paths = [
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
]

for tesseract_path in possible_tesseract_paths:
    if os.path.exists(tesseract_path):
        pytesseract.pytesseract.tesseract_cmd = tesseract_path
        break


# --------------------------------------------------
# OCR TEXT EXTRACTION
# --------------------------------------------------


def extract_ocr_rows(page):
    matrix = pymupdf.Matrix(2.5, 2.5)

    pixmap = page.get_pixmap(matrix=matrix, alpha=False)

    image_bytes = pixmap.tobytes("png")

    image = PILImage.open(BytesIO(image_bytes))

    ocr_data = pytesseract.image_to_data(
        image, output_type=pytesseract.Output.DICT, config="--psm 6"
    )

    lines = {}

    total_words = len(ocr_data["text"])

    for index in range(total_words):

        text = ocr_data["text"][index].strip()

        if not text:
            continue

        try:
            confidence = float(ocr_data["conf"][index])
        except Exception:
            confidence = 0

        if confidence < 20:
            continue

        key = (
            ocr_data["block_num"][index],
            ocr_data["par_num"][index],
            ocr_data["line_num"][index],
        )

        word_data = {
            "text": text,
            "left": ocr_data["left"][index],
            "width": ocr_data["width"][index],
        }

        lines.setdefault(key, []).append(word_data)

    rows = []

    for key in sorted(lines.keys()):

        words = sorted(lines[key], key=lambda item: item["left"])

        if not words:
            continue

        row = []
        current_column = []
        previous_right = None

        for word in words:

            left = word["left"]
            right = left + word["width"]

            if previous_right is not None:
                gap = left - previous_right

                if gap > 90:
                    if current_column:
                        row.append(" ".join(item["text"] for item in current_column))

                    current_column = []

            current_column.append(word)

            previous_right = right

        if current_column:
            row.append(" ".join(item["text"] for item in current_column))

        if row:
            rows.append(row)

    return rows


# --------------------------------------------------
# EXCEL FORMATTING
# --------------------------------------------------


def format_excel_sheet(worksheet):

    for row in worksheet.iter_rows():

        for cell in row:

            cell.alignment = Alignment(vertical="top", wrap_text=True)

    for column_number in range(1, worksheet.max_column + 1):

        column_letter = get_column_letter(column_number)

        max_length = 0

        for cell in worksheet[column_letter]:

            if cell.value is not None:

                value_length = len(str(cell.value))

                if value_length > max_length:
                    max_length = value_length

        worksheet.column_dimensions[column_letter].width = min(
            max(max_length + 3, 12), 45
        )


# --------------------------------------------------
# ADD ORIGINAL PAGE PREVIEW
# --------------------------------------------------


def add_page_preview(worksheet, page, image_buffers):

    matrix = pymupdf.Matrix(1.5, 1.5)

    pixmap = page.get_pixmap(matrix=matrix, alpha=False)

    preview_bytes = pixmap.tobytes("png")

    preview_buffer = BytesIO(preview_bytes)

    image_buffers.append(preview_buffer)

    excel_image = XLImage(preview_buffer)

    max_width = 420
    max_height = 600

    scale = min(max_width / excel_image.width, max_height / excel_image.height, 1)

    excel_image.width = int(excel_image.width * scale)

    excel_image.height = int(excel_image.height * scale)

    worksheet["E1"] = "Original PDF Page Preview"

    worksheet["E1"].font = Font(bold=True)

    worksheet.add_image(excel_image, "E2")


# --------------------------------------------------
# ADD EMBEDDED PDF IMAGES / LOGOS
# --------------------------------------------------


def add_embedded_images(worksheet, document, page, image_buffers):

    images = page.get_images(full=True)

    if not images:
        return

    start_row = 35

    worksheet[f"E{start_row}"] = "Extracted Images / Logos"

    worksheet[f"E{start_row}"].font = Font(bold=True)

    current_row = start_row + 2

    used_xrefs = set()

    for image_info in images:

        xref = image_info[0]

        if xref in used_xrefs:
            continue

        used_xrefs.add(xref)

        try:

            extracted = document.extract_image(xref)

            image_bytes = extracted["image"]

            pil_image = PILImage.open(BytesIO(image_bytes))

            if pil_image.width < 40 or pil_image.height < 40:
                continue

            converted_buffer = BytesIO()

            if pil_image.mode not in ("RGB", "RGBA"):
                pil_image = pil_image.convert("RGB")

            pil_image.save(converted_buffer, format="PNG")

            converted_buffer.seek(0)

            image_buffers.append(converted_buffer)

            excel_image = XLImage(converted_buffer)

            max_width = 250
            max_height = 200

            scale = min(
                max_width / excel_image.width, max_height / excel_image.height, 1
            )

            excel_image.width = int(excel_image.width * scale)

            excel_image.height = int(excel_image.height * scale)

            anchor = f"E{current_row}"

            worksheet.add_image(excel_image, anchor)

            rows_needed = max(10, int(excel_image.height / 20) + 2)

            current_row += rows_needed

        except Exception as image_error:

            print("PDF image extraction error:", image_error)


# --------------------------------------------------
# PDF TO EXCEL ENDPOINT
# --------------------------------------------------


@app.post("/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):

    document = None

    try:

        if not file.filename:

            raise HTTPException(status_code=400, detail="Please select a PDF file.")

        if not file.filename.lower().endswith(".pdf"):

            raise HTTPException(
                status_code=400, detail="Please select a valid PDF file."
            )

        pdf_bytes = await file.read()

        if not pdf_bytes:

            raise HTTPException(
                status_code=400, detail="The selected PDF file is empty."
            )

        document = pymupdf.open(stream=pdf_bytes, filetype="pdf")

        if document.needs_pass:

            raise HTTPException(
                status_code=400,
                detail=("This PDF is password protected. " "Please unlock it first."),
            )

        if document.page_count == 0:

            raise HTTPException(status_code=400, detail="This PDF has no pages.")

        workbook = Workbook()

        default_sheet = workbook.active

        workbook.remove(default_sheet)

        image_buffers = []

        total_content_found = False

        # --------------------------------------------------
        # PROCESS EVERY PDF PAGE
        # --------------------------------------------------

        for page_number in range(document.page_count):

            page = document[page_number]

            worksheet = workbook.create_sheet(title=f"Page {page_number + 1}")

            content_found = False

            # ----------------------------------------------
            # STEP 1 - TRY TABLE EXTRACTION
            # ----------------------------------------------

            try:

                tables = page.find_tables()

                table_number = 1

                for table in tables.tables:

                    extracted_table = table.extract()

                    if not extracted_table:
                        continue

                    if table_number > 1:
                        worksheet.append([])

                    worksheet.append([f"Table {table_number}"])

                    worksheet.cell(row=worksheet.max_row, column=1).font = Font(
                        bold=True
                    )

                    for table_row in extracted_table:

                        cleaned_row = []

                        for value in table_row:

                            if value is None:
                                cleaned_row.append("")
                            else:
                                cleaned_row.append(str(value).strip())

                        worksheet.append(cleaned_row)

                    table_number += 1

                    content_found = True

            except Exception as table_error:

                print("PDF table extraction error:", table_error)

            # ----------------------------------------------
            # STEP 2 - NORMAL PDF TEXT
            # ----------------------------------------------

            normal_text = page.get_text("text").strip()

            if normal_text:

                if content_found:
                    worksheet.append([])
                    worksheet.append([])

                worksheet.append(["Extracted Text"])

                worksheet.cell(row=worksheet.max_row, column=1).font = Font(bold=True)

                for text_line in normal_text.splitlines():

                    text_line = text_line.strip()

                    if text_line:

                        worksheet.append([text_line])

                content_found = True

            # ----------------------------------------------
            # STEP 3 - OCR FOR SCANNED / IMAGE PDF
            # ----------------------------------------------

            if not normal_text:

                try:

                    ocr_rows = extract_ocr_rows(page)

                    if ocr_rows:

                        worksheet.append(["OCR Extracted Data"])

                        worksheet.cell(row=worksheet.max_row, column=1).font = Font(
                            bold=True
                        )

                        for ocr_row in ocr_rows:

                            worksheet.append(ocr_row)

                        content_found = True

                except pytesseract.TesseractNotFoundError:

                    raise HTTPException(
                        status_code=500,
                        detail=(
                            "Tesseract OCR is not installed " "or could not be found."
                        ),
                    )

                except Exception as ocr_error:

                    print("OCR extraction error:", ocr_error)

            # ----------------------------------------------
            # STEP 4 - ORIGINAL PAGE VISUAL
            # ----------------------------------------------

            try:

                add_page_preview(worksheet, page, image_buffers)

                content_found = True

            except Exception as preview_error:

                print("PDF page preview error:", preview_error)

            # ----------------------------------------------
            # STEP 5 - EXTRACT PHOTO / LOGO / IMAGES
            # ----------------------------------------------

            try:

                add_embedded_images(worksheet, document, page, image_buffers)

            except Exception as embedded_error:

                print("Embedded image error:", embedded_error)

            # ----------------------------------------------
            # EMPTY PAGE
            # ----------------------------------------------

            if not content_found:

                worksheet.append(
                    ["No readable text, table, " "or image was detected on this page."]
                )

            format_excel_sheet(worksheet)

            total_content_found = total_content_found or content_found

        # --------------------------------------------------
        # SAVE EXCEL FILE
        # --------------------------------------------------

        if not total_content_found:

            raise HTTPException(
                status_code=400,
                detail=("No readable content could " "be detected in this PDF."),
            )

        output = BytesIO()

        workbook.save(output)

        output.seek(0)

        original_name = file.filename.rsplit(".", 1)[0]

        safe_name = "".join(
            character if character.isalnum() or character in ("-", "_", " ") else "_"
            for character in original_name
        ).strip()

        if not safe_name:
            safe_name = "converted"

        output_filename = f"{safe_name}.xlsx"

        return StreamingResponse(
            output,
            media_type=(
                "application/"
                "vnd.openxmlformats-officedocument."
                "spreadsheetml.sheet"
            ),
            headers={
                "Content-Disposition": f'attachment; filename="{output_filename}"'
            },
        )

    except HTTPException:

        raise

    except pytesseract.TesseractNotFoundError:

        raise HTTPException(
            status_code=500,
            detail=("Tesseract OCR is not installed " "or could not be found."),
        )

    except Exception as error:

        print("PDF to Excel Error:", error)

        raise HTTPException(
            status_code=500, detail="Unable to convert this PDF to Excel."
        )

    finally:

        if document is not None:

            try:
                document.close()

            except Exception:
                pass
